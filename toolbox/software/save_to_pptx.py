# author: yannik fontana, creation date: 06.05.2026
"""
Append Plotly figures to the daily lab-journal PowerPoint (``.pptx``).

**Dependencies** (install in your environment):

- ``python-pptx`` — build/save presentations.
- ``plotly`` — figure objects.
- ``kaleido`` — static image export (``Figure.write_image(..., format="png")``).
  Without Kaleido, PNG export raises at runtime; install with ``pip install kaleido``.

``pptx`` / ``plotly`` are imported only when :func:`save_to_pptx` runs, so
:func:`labjournal_pptx_path` works without them.

Layout uses absolute positions in centimetres; tune the ``*_CM`` constants below.
"""
from __future__ import annotations

import datetime
import io
import logging
from pathlib import Path
from typing import Any, Sequence

from toolbox.software.path_config import get_labjournalpath

logger = logging.getLogger(__name__)
TITLE_LEFT_CM = 0.8
TITLE_TOP_CM = 0.6
TITLE_HEIGHT_CM = 1.2

TAG_MARGIN_RIGHT_CM = 0.8
TAG_TOP_CM = 1.5
TAG_BOX_WIDTH_CM = 7.0
TAG_BOX_HEIGHT_CM = 4.5
TAG_FONT_PT = 10

FIGURE_LEFT_CM = 0.8
FIGURE_TOP_CM = 2.2
FIGURE_WIDTH_CM = 20.0


def labjournal_pptx_path(journal_date: datetime.date | None = None) -> Path:
    """
    ``<labjournalpath>/<YYYY>/<YYYYMMDD>.pptx`` (local calendar date, naive).
    """
    d = journal_date or datetime.date.today()
    root = Path(get_labjournalpath())
    year_dir = root / f"{d.year:04d}"
    year_dir.mkdir(parents=True, exist_ok=True)
    return year_dir / f"{d:%Y%m%d}.pptx"


def _slide_title_line(data: Any) -> str:
    ts = getattr(data, "timestamp", None)
    if isinstance(ts, datetime.datetime):
        date_s = ts.strftime("%d.%m.%Y")
    else:
        date_s = datetime.date.today().strftime("%d.%m.%Y")
    name = getattr(data, "filename", "") or ""
    return f"{date_s} {name}".strip()


def _tag_text(data: Any) -> str:
    tag = getattr(data, "tag", None)
    if not tag:
        return ""
    if isinstance(tag, str):
        return tag
    if isinstance(tag, Sequence) and not isinstance(tag, (bytes, str)):
        lines = [str(x) for x in tag]
        return "\n".join(lines)
    return str(tag)


def save_to_pptx(
    data: Any,
    figures: Sequence[Any],
    *,
    journal_date: datetime.date | None = None,
) -> Path | None:
    """
    For each Plotly figure, append a blank slide with title
    ``DD.MM.YYYY <filename>``, tag (upper right), and the rasterised figure.

    If ``figures`` is empty, logs a warning and returns ``None`` without touching
    the journal file.

    Returns the path to the saved ``.pptx``, or ``None`` when no figures were given.

    Raises:
        ImportError: if ``python-pptx`` or ``plotly`` is not installed.
        TypeError: if an element of ``figures`` is not a ``plotly.graph_objects.Figure``.
    """
    if not figures:
        logger.warning(
            "No figures provided; did not write anything to the lab journal."
        )
        return None

    try:
        from pptx import Presentation
        from pptx.util import Cm, Pt
        from plotly.graph_objects import Figure as GoFigure
    except ImportError as exc:
        raise ImportError(
            "save_to_pptx requires python-pptx and plotly "
            "(and kaleido for PNG export). "
            "Install e.g.: pip install python-pptx plotly kaleido"
        ) from exc

    def plotly_to_png(fig: Any) -> io.BytesIO:
        buf = io.BytesIO()
        fig.write_image(buf, format="png", scale=2)
        buf.seek(0)
        return buf

    def add_figure_slide(prs: Any, fig: Any) -> None:
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)
        sw = int(prs.slide_width)

        title_left = Cm(TITLE_LEFT_CM)
        title_top = Cm(TITLE_TOP_CM)
        tag_left_edge = sw - int(Cm(TAG_MARGIN_RIGHT_CM + TAG_BOX_WIDTH_CM))
        margin = int(Cm(0.5))
        title_w = max(tag_left_edge - int(title_left) - margin, int(Cm(5)))
        title_box = slide.shapes.add_textbox(
            title_left, title_top, title_w, Cm(TITLE_HEIGHT_CM)
        )
        title_box.text_frame.text = _slide_title_line(data)
        title_box.text_frame.word_wrap = True

        tag_left = sw - int(Cm(TAG_MARGIN_RIGHT_CM + TAG_BOX_WIDTH_CM))
        tag_top = Cm(TAG_TOP_CM)
        tag_box = slide.shapes.add_textbox(
            tag_left, tag_top, Cm(TAG_BOX_WIDTH_CM), Cm(TAG_BOX_HEIGHT_CM)
        )
        tag_box.text_frame.text = _tag_text(data) or " "
        tag_box.text_frame.word_wrap = True
        for p in tag_box.text_frame.paragraphs:
            p.font.size = Pt(TAG_FONT_PT)

        png = plotly_to_png(fig)
        slide.shapes.add_picture(
            png, Cm(FIGURE_LEFT_CM), Cm(FIGURE_TOP_CM), width=Cm(FIGURE_WIDTH_CM)
        )

    path = labjournal_pptx_path(journal_date)
    prs = Presentation(str(path)) if path.is_file() else Presentation()

    for fig in figures:
        if not isinstance(fig, GoFigure):
            raise TypeError(
                f"expected plotly.graph_objects.Figure, got {type(fig).__name__}"
            )
        add_figure_slide(prs, fig)

    prs.save(str(path))
    logger.info("Lab journal updated: %s (%d slide(s))", path, len(figures))
    return path
